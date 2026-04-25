#!/usr/bin/env python3
"""Safely modify Office documents by writing a new output file."""

import argparse
import json
import os
import re


def ensure_dir(path):
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)


def read_spec(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


def assert_safe_output(input_path, output_path):
    if os.path.abspath(input_path) == os.path.abspath(output_path):
        raise SystemExit('Refusing to overwrite the input file; choose a different output path.')
    if os.path.exists(output_path):
        raise SystemExit(f'Refusing to overwrite existing output: {output_path}')


def replace_in_paragraph(paragraph, old, new):
    if old not in paragraph.text:
        return False
    # Rebuild runs to keep behavior predictable for generated documents.
    paragraph.text = paragraph.text.replace(old, new)
    return True


def add_toc(document, title='目录'):
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    document.add_heading(title, level=1)
    paragraph = document.add_paragraph()

    begin = OxmlElement('w:fldChar')
    begin.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText')
    instr.set(qn('xml:space'), 'preserve')
    instr.text = 'TOC \\o "1-3" \\h \\z \\u'
    separate = OxmlElement('w:fldChar')
    separate.set(qn('w:fldCharType'), 'separate')
    end = OxmlElement('w:fldChar')
    end.set(qn('w:fldCharType'), 'end')

    paragraph._p.append(begin)
    paragraph._p.append(instr)
    paragraph._p.append(separate)
    run = paragraph.add_run('目录字段已插入；在 Word 中打开后可右键更新域。')
    run.italic = True
    paragraph._p.append(end)


def apply_docx_heading_style(document, op):
    from docx.shared import RGBColor

    color = str(op.get('color', '1F4E79')).lstrip('#')
    if not re.fullmatch(r'[0-9A-Fa-f]{6}', color):
        color = '1F4E79'
    font_name = op.get('font')
    for style_name in ['Heading 1', 'Heading 2', 'Title']:
        try:
            style = document.styles[style_name]
        except KeyError:
            continue
        style.font.color.rgb = RGBColor.from_string(color.upper())
        if font_name:
            style.font.name = str(font_name)


def modify_docx(input_path, spec, output_path):
    from docx import Document

    document = Document(input_path)
    for op in spec.get('operations', []):
        name = op.get('op')
        if name == 'replace_text':
            old = str(op.get('old', ''))
            new = str(op.get('new', ''))
            if not old:
                continue
            for paragraph in document.paragraphs:
                replace_in_paragraph(paragraph, old, new)
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            replace_in_paragraph(paragraph, old, new)
        elif name == 'append_section':
            document.add_heading(op.get('heading', 'Section'), level=int(op.get('level', 1)))
            for para in op.get('paragraphs', []):
                document.add_paragraph(str(para))
            for item in op.get('bullets', []):
                document.add_paragraph(str(item), style='List Bullet')
            table = op.get('table')
            if table:
                headers = table.get('headers', [])
                rows = table.get('rows', [])
                cols = max(len(headers), *(len(row) for row in rows), 1)
                doc_table = document.add_table(rows=1 if headers else 0, cols=cols)
                doc_table.style = 'Table Grid'
                if headers:
                    for i, value in enumerate(headers):
                        doc_table.rows[0].cells[i].text = str(value)
                for row in rows:
                    cells = doc_table.add_row().cells
                    for i, value in enumerate(row):
                        cells[i].text = str(value)
        elif name == 'add_toc':
            add_toc(document, op.get('title', '目录'))
        elif name == 'set_heading_style':
            apply_docx_heading_style(document, op)
    ensure_dir(output_path)
    document.save(output_path)


def get_or_create_sheet(workbook, name):
    if name in workbook.sheetnames:
        return workbook[name]
    return workbook.create_sheet(title=str(name)[:31])


def apply_header_style(sheet, op):
    from openpyxl.styles import Font, PatternFill

    fill = str(op.get('fill', '1F4E79')).lstrip('#')
    font_color = str(op.get('font_color', 'FFFFFF')).lstrip('#')
    if not re.fullmatch(r'[0-9A-Fa-f]{6}', fill):
        fill = '1F4E79'
    if not re.fullmatch(r'[0-9A-Fa-f]{6}', font_color):
        font_color = 'FFFFFF'
    for cell in sheet[1]:
        cell.font = Font(bold=True, color=font_color.upper())
        cell.fill = PatternFill('solid', fgColor=fill.upper())


def modify_xlsx(input_path, spec, output_path):
    from openpyxl import load_workbook

    workbook = load_workbook(input_path)
    for op in spec.get('operations', []):
        name = op.get('op')
        sheet_name = op.get('sheet') or workbook.sheetnames[0]
        sheet = get_or_create_sheet(workbook, sheet_name)
        if name == 'set_cell':
            sheet[str(op.get('cell', 'A1'))] = op.get('value')
        elif name == 'insert_formula':
            formula = str(op.get('formula', '')).strip()
            if formula and not formula.startswith('='):
                formula = '=' + formula
            sheet[str(op.get('cell', 'A1'))] = formula
        elif name == 'append_rows':
            for row in op.get('rows', []):
                sheet.append(row)
        elif name == 'style_header':
            apply_header_style(sheet, op)
        elif name == 'replace_text':
            old = str(op.get('old', ''))
            new = str(op.get('new', ''))
            if not old:
                continue
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and old in cell.value:
                        cell.value = cell.value.replace(old, new)
    ensure_dir(output_path)
    workbook.save(output_path)


def hex_to_rgb(hex_color, fallback='1F4E79'):
    value = str(hex_color or fallback).lstrip('#')
    if not re.fullmatch(r'[0-9A-Fa-f]{6}', value):
        value = fallback
    return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))


def apply_pptx_theme(prs, op):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    primary = RGBColor(*hex_to_rgb(op.get('primary_color'), '1F4E79'))
    accent = RGBColor(*hex_to_rgb(op.get('accent_color'), 'F28C28'))
    for slide in prs.slides:
        if slide.shapes.title:
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = primary
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()


def modify_pptx(input_path, spec, output_path):
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(input_path)
    for op in spec.get('operations', []):
        name = op.get('op')
        if name == 'replace_text':
            old = str(op.get('old', ''))
            new = str(op.get('new', ''))
            if not old:
                continue
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and old in shape.text:
                        shape.text = shape.text.replace(old, new)
        elif name == 'append_slide':
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = str(op.get('title', 'Untitled'))
            body = slide.placeholders[1].text_frame
            body.clear()
            body_text = op.get('body')
            if body_text:
                body.paragraphs[0].text = str(body_text)
                body.paragraphs[0].font.size = Pt(20)
            for item in op.get('bullets', []):
                p = body.add_paragraph()
                p.text = str(item)
                p.level = 0
                p.font.size = Pt(20)
            notes = op.get('notes')
            if notes:
                slide.notes_slide.notes_text_frame.text = str(notes)
        elif name == 'set_theme':
            apply_pptx_theme(prs, op)
    ensure_dir(output_path)
    prs.save(output_path)


def main():
    parser = argparse.ArgumentParser(description='Modify Office documents without overwriting originals.')
    parser.add_argument('input')
    parser.add_argument('spec')
    parser.add_argument('output')
    args = parser.parse_args()
    assert_safe_output(args.input, args.output)
    spec = read_spec(args.spec)
    ext = os.path.splitext(args.input)[1].lower()
    if ext == '.docx':
        modify_docx(args.input, spec, args.output)
    elif ext == '.xlsx':
        modify_xlsx(args.input, spec, args.output)
    elif ext == '.pptx':
        modify_pptx(args.input, spec, args.output)
    else:
        raise SystemExit(f'Unsupported input extension: {ext}')
    print(args.output)


if __name__ == '__main__':
    main()
