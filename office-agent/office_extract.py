#!/usr/bin/env python3
"""Extract text from generated Office or fallback HTML documents."""

import argparse
import json
import os
from html.parser import HTMLParser


class TextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data):
        text = data.strip()
        if text:
            self.parts.append(text)


def extract_html(path):
    parser_obj = TextExtractor()
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        parser_obj.feed(f.read())
    return parser_obj.parts


def extract_docx(path):
    from docx import Document

    document = Document(path)
    parts = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                parts.append(' | '.join(values))
    return parts


def extract_xlsx(path):
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                parts.append(' | '.join(values))
    return parts


def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"[Slide {index}]")
            parts.extend(slide_parts)
    return parts


def extract(path):
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == '.docx':
            return extract_docx(path)
        if ext == '.xlsx':
            return extract_xlsx(path)
        if ext == '.pptx':
            return extract_pptx(path)
    except Exception:
        return extract_html(path)
    return extract_html(path)


def main():
    parser = argparse.ArgumentParser(description='Extract text outline from generated Office files.')
    parser.add_argument('input')
    parser.add_argument('output')
    args = parser.parse_args()
    result = {'source': args.input, 'text': extract(args.input)}
    with open(args.output, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(args.output)


if __name__ == '__main__':
    main()
