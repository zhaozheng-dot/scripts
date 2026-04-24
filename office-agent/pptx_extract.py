#!/usr/bin/env python3
"""Extract PPTX content and source map for Office Agent conversions."""

import argparse
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
if SCRIPT_DIR not in sys.path:
    sys.path.insert(0, SCRIPT_DIR)

from office_common import ensure_dir, write_json


def guess_semantic(text):
    low = text.lower()
    if any(key in low for key in ['risk', '风险', 'pending', '警示']):
        return 'risk'
    if any(key in low for key in ['source', '来源', '置信度', 'disclaimer', '免责声明']):
        return 'source_or_disclaimer'
    if any(key in low for key in ['summary', '摘要', '结论', '建议']):
        return 'summary_or_recommendation'
    if any(key in low for key in ['swot', 'strength', 'weakness', 'opportunit', 'threat']):
        return 'swot'
    if len(text) <= 32:
        return 'label_or_heading'
    return 'body'


def bbox(shape):
    return [int(shape.left), int(shape.top), int(shape.width), int(shape.height)]


def text_items(shape):
    items = []
    if not getattr(shape, 'has_text_frame', False):
        return items
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        items.append({
            'type': 'text',
            'text': text,
            'paragraph_index': paragraph_index,
            'level': paragraph.level,
            'bbox': bbox(shape),
            'semantic_guess': guess_semantic(text),
            'confidence': 0.7,
        })
    return items


def table_item(shape):
    if not getattr(shape, 'has_table', False):
        return None
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return {'type': 'table', 'rows': rows, 'bbox': bbox(shape), 'semantic_guess': 'table', 'confidence': 0.9}


def save_picture(shape, assets_dir, slide_no, image_index):
    image = shape.image
    ext = image.ext or 'png'
    name = f'slide-{slide_no:03d}-image-{image_index:02d}.{ext}'
    path = os.path.join(assets_dir, name)
    ensure_dir(path)
    with open(path, 'wb') as f:
        f.write(image.blob)
    return {'type': 'image', 'path': path, 'ext': ext, 'bbox': bbox(shape), 'semantic_guess': 'image', 'confidence': 0.9}


def notes_text(slide):
    try:
        text = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ''
    return text


def extract(input_path, assets_dir=None):
    prs = Presentation(input_path)
    slides = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title is not None else ''
        items = []
        image_index = 1
        for shape_index, shape in enumerate(sorted(slide.shapes, key=lambda s: (getattr(s, 'top', 0), getattr(s, 'left', 0))), start=1):
            for item in text_items(shape):
                item['shape_index'] = shape_index
                items.append(item)
            table = table_item(shape)
            if table:
                table['shape_index'] = shape_index
                items.append(table)
            if assets_dir and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = save_picture(shape, assets_dir, slide_no, image_index)
                image['shape_index'] = shape_index
                items.append(image)
                image_index += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                items.append({'type': 'image', 'bbox': bbox(shape), 'semantic_guess': 'image', 'confidence': 0.8, 'shape_index': shape_index})
            elif shape.shape_type in {MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.CHART}:
                items.append({'type': str(shape.shape_type), 'bbox': bbox(shape), 'semantic_guess': 'complex_visual', 'confidence': 0.5, 'shape_index': shape_index})
        slides.append({'slide_no': slide_no, 'title': title, 'notes': notes_text(slide), 'items': items})
    return {'source': input_path, 'slides': slides}


def main():
    parser = argparse.ArgumentParser(description='Extract PPTX content and source map.')
    parser.add_argument('input')
    parser.add_argument('output')
    parser.add_argument('--assets-dir', default='')
    args = parser.parse_args()
    data = extract(args.input, args.assets_dir or None)
    write_json(args.output, data)
    print(args.output)


if __name__ == '__main__':
    main()
