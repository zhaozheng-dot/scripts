#!/usr/bin/env python3
"""Create synthetic but business-like Office files for service regression."""

import argparse
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DEFAULT_OUTPUT_DIR = os.path.join(SCRIPT_DIR, 'examples', 'service_regression_inputs')


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def build_project_weekly_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = '脱敏项目周报'
    slide.placeholders[1].text = 'P2 服务化回归样例 / Synthetic Business Fixture'

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '本周进展'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for index, item in enumerate(['完成 HTTP 服务可靠性增强', '建立任务事件审计日志', '生成/转换/修改链路进入回归验证', '等待真实业务材料补充']):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = '风险台账摘要'
    table = slide.shapes.add_table(4, 4, Inches(0.6), Inches(1.4), Inches(8.4), Inches(2.0)).table
    rows = [
        ['风险', '等级', '负责人', '缓解措施'],
        ['样例不足', '中', 'Owner A', '使用合成样例补齐基线'],
        ['复杂版式退化', '中', 'Owner B', '保留 ledger 并人工复核'],
        ['真实数据敏感', '高', 'Owner C', '仅接收脱敏材料'],
    ]
    for r, values in enumerate(rows):
        for c, value in enumerate(values):
            table.cell(r, c).text = value

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '下周计划'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for index, item in enumerate(['扩展异常回归', '沉淀模板化生成场景', '对失败样例归类形成 P3/P4 路线图']):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
    prs.save(path)


def build_decision_review_pptx(path):
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = '脱敏决策评审材料'
    slide.placeholders[1].text = '高密度文本 + 指标图表 / Synthetic Fixture'

    for idx in range(1, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f'评审要点 {idx}'
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for j in range(1, 9):
            p = tf.paragraphs[0] if j == 1 else tf.add_paragraph()
            p.text = f'脱敏评审条目 {idx}.{j}: 目标、收益、成本、风险、依赖、复核口径均需记录。'
            p.level = 0

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = '指标趋势'
    data = CategoryChartData()
    data.categories = ['M1', 'M2', 'M3', 'M4']
    data.add_series('计划', (80, 85, 90, 95))
    data.add_series('实际', (76, 82, 88, 91))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(7.2), Inches(4.2), data)
    prs.save(path)


def build_meeting_minutes_docx(path):
    from docx import Document

    doc = Document()
    doc.add_heading('脱敏项目会议纪要', level=0)
    doc.add_paragraph('会议主题：Office Agent 服务化回归建设')
    doc.add_paragraph('会议日期：2026-04-25')
    doc.add_heading('结论摘要', level=1)
    for item in ['优先建立合成真实感样例', '服务层需持续输出 events.jsonl', '失败样例按能力短板归类']:
        doc.add_paragraph(item, style='List Bullet')
    doc.add_heading('行动项', level=1)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    for i, value in enumerate(['事项', '负责人', '截止时间', '状态']):
        table.rows[0].cells[i].text = value
    for row in [
        ['新增服务回归脚本', 'Owner A', '本周', '待确认'],
        ['补齐异常用例', 'Owner B', '下周', '待确认'],
        ['沉淀模板路线图', 'Owner C', '下周', '进行中'],
    ]:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = value
    doc.save(path)


def build_project_plan_docx(path):
    from docx import Document

    doc = Document()
    doc.add_heading('脱敏项目实施方案', level=0)
    for heading, paragraphs in [
        ('背景', ['当前系统已具备 CLI、HTTP、MCP-style 三种入口，需要通过回归样例提升稳定性。']),
        ('目标', ['建立可重复执行的服务级回归套件。', '覆盖生成、转换、修改、取消和异常处理。']),
        ('范围', ['仅使用合成或脱敏材料。', '不提交任何真实客户或敏感文件。']),
        ('验收标准', ['服务回归状态为 pass。', '每个成功任务包含完整事件链。', '失败任务输出结构化错误。']),
    ]:
        doc.add_heading(heading, level=1)
        for paragraph in paragraphs:
            doc.add_paragraph(paragraph)
    doc.save(path)


def build_budget_tracker_xlsx(path):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = '预算'
    ws.append(['科目', '预算', '实际', '差异', '状态'])
    for row in [
        ['人力', 120000, 98000, 22000, '正常'],
        ['工具', 30000, 31500, -1500, '待确认'],
        ['外包', 50000, 26000, 24000, '正常'],
        ['培训', 15000, 18000, -3000, '关注'],
    ]:
        ws.append(row)
    for cell in ws[1]:
        cell.font = Font(bold=True, color='FFFFFF')
        cell.fill = PatternFill('solid', fgColor='1F4E79')
    risk = wb.create_sheet('风险')
    risk.append(['编号', '风险', '等级', '状态'])
    risk.append([1, '预算超支', '中', '待确认'])
    risk.append([2, '采购延期', '低', '监控中'])
    wb.save(path)


def build_risk_register_xlsx(path):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = '风险台账'
    ws.append(['编号', '类别', '描述', '等级', '责任人', '状态'])
    rows = [
        [1, '数据', '真实样例不足导致覆盖不足', '中', 'Owner A', '待确认'],
        [2, '质量', '复杂图表转换可能丢失细节', '中', 'Owner B', '监控中'],
        [3, '流程', '用户确认门控被绕过', '高', 'Owner C', '已控制'],
        [4, '集成', '上层 Agent 轮询超时', '中', 'Owner D', '待确认'],
    ]
    for row in rows:
        ws.append(row)
    for cell in ws[1]:
        cell.font = Font(bold=True, color='FFFFFF')
        cell.fill = PatternFill('solid', fgColor='7A1F1F')
    wb.save(path)


def main(argv=None):
    parser = argparse.ArgumentParser(description='Create synthetic service regression Office fixtures.')
    parser.add_argument('--output-dir', default=DEFAULT_OUTPUT_DIR)
    args = parser.parse_args(argv)
    ensure_dir(args.output_dir)
    fixtures = {
        'business_project_weekly.pptx': build_project_weekly_pptx,
        'business_decision_review.pptx': build_decision_review_pptx,
        'business_meeting_minutes.docx': build_meeting_minutes_docx,
        'business_project_plan.docx': build_project_plan_docx,
        'business_budget_tracker.xlsx': build_budget_tracker_xlsx,
        'business_risk_register.xlsx': build_risk_register_xlsx,
    }
    for name, builder in fixtures.items():
        builder(os.path.join(args.output_dir, name))
    print(args.output_dir)


if __name__ == '__main__':
    main()
