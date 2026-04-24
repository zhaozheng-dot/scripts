#!/usr/bin/env python3
"""Generate Office documents from JSON specs."""

import argparse
import html
import json
import os
from datetime import datetime


def ensure_dir(path):
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)


def esc(value):
    return html.escape(str(value), quote=True)


def read_spec(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


def html_page(title, body, extra_style=''):
    style = f"""
body {{ font-family: Aptos, Calibri, Arial, sans-serif; line-height: 1.55; color: #1f2933; margin: 48px; }}
h1 {{ color: #12355b; border-bottom: 2px solid #d7e3f4; padding-bottom: 8px; }}
h2 {{ color: #1f4f7a; margin-top: 28px; }}
table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
th, td {{ border: 1px solid #b8c6d8; padding: 8px 10px; vertical-align: top; }}
th {{ background: #eaf2fb; }}
.slide {{ page-break-after: always; min-height: 560px; padding: 28px; border: 1px solid #d5dee8; margin-bottom: 28px; }}
.subtitle {{ color: #64748b; font-size: 18px; }}
.meta {{ color: #64748b; font-size: 12px; }}
{extra_style}
"""
    return f"<!doctype html><html><head><meta charset='utf-8'><title>{esc(title)}</title><style>{style}</style></head><body>{body}</body></html>"


def render_docx_html(spec, output):
    title = spec.get('title', 'Document')
    parts = [f"<h1>{esc(title)}</h1>", f"<p class='meta'>Generated: {datetime.now():%Y-%m-%d %H:%M}</p>"]
    for section in spec.get('sections', []):
        parts.append(f"<h2>{esc(section.get('heading', 'Section'))}</h2>")
        for para in section.get('paragraphs', []):
            parts.append(f"<p>{esc(para)}</p>")
        table = section.get('table')
        if table:
            parts.append('<table>')
            headers = table.get('headers', [])
            if headers:
                parts.append('<tr>' + ''.join(f"<th>{esc(h)}</th>" for h in headers) + '</tr>')
            for row in table.get('rows', []):
                parts.append('<tr>' + ''.join(f"<td>{esc(c)}</td>" for c in row) + '</tr>')
            parts.append('</table>')
        for item in section.get('bullets', []):
            parts.append(f"<ul><li>{esc(item)}</li></ul>")
    ensure_dir(output)
    with open(output, 'w', encoding='utf-8') as f:
        f.write(html_page(title, '\n'.join(parts)))


def render_xlsx_html(spec, output):
    title = spec.get('title', 'Workbook')
    parts = [f"<h1>{esc(title)}</h1>"]
    for sheet in spec.get('sheets', []):
        parts.append(f"<h2>{esc(sheet.get('name', 'Sheet'))}</h2><table>")
        headers = sheet.get('headers', [])
        if headers:
            parts.append('<tr>' + ''.join(f"<th>{esc(h)}</th>" for h in headers) + '</tr>')
        for row in sheet.get('rows', []):
            parts.append('<tr>' + ''.join(f"<td>{esc(c)}</td>" for c in row) + '</tr>')
        parts.append('</table>')
    ensure_dir(output)
    with open(output, 'w', encoding='utf-8') as f:
        f.write(html_page(title, '\n'.join(parts)))


def render_pptx_html(spec, output):
    title = spec.get('title', 'Slide Deck')
    parts = [f"<h1>{esc(title)}</h1>"]
    for index, slide in enumerate(spec.get('slides', []), start=1):
        parts.append('<section class="slide">')
        parts.append(f"<p class='meta'>Slide {index}</p>")
        parts.append(f"<h2>{esc(slide.get('title', 'Untitled'))}</h2>")
        if slide.get('subtitle'):
            parts.append(f"<p class='subtitle'>{esc(slide['subtitle'])}</p>")
        if slide.get('body'):
            parts.append(f"<p>{esc(slide['body'])}</p>")
        bullets = slide.get('bullets', [])
        if bullets:
            parts.append('<ul>' + ''.join(f"<li>{esc(b)}</li>" for b in bullets) + '</ul>')
        parts.append('</section>')
    ensure_dir(output)
    with open(output, 'w', encoding='utf-8') as f:
        f.write(html_page(title, '\n'.join(parts), extra_style='@page { size: landscape; }'))


def render_docx_ooxml(spec, output):
    from docx import Document
    from docx.shared import Inches

    document = Document()
    title = spec.get('title', 'Document')
    document.add_heading(title, level=0)
    document.add_paragraph(f"Generated: {datetime.now():%Y-%m-%d %H:%M}")
    for section in spec.get('sections', []):
        document.add_heading(section.get('heading', 'Section'), level=1)
        for para in section.get('paragraphs', []):
            document.add_paragraph(str(para))
        for item in section.get('bullets', []):
            document.add_paragraph(str(item), style='List Bullet')
        table = section.get('table')
        if table:
            headers = table.get('headers', [])
            rows = table.get('rows', [])
            cols = max(len(headers), *(len(row) for row in rows), 1)
            doc_table = document.add_table(rows=1 if headers else 0, cols=cols)
            doc_table.style = 'Table Grid'
            if headers:
                for i, value in enumerate(headers):
                    doc_table.rows[0].cells[i].text = str(value)
            for row in rows:
                cells = doc_table.add_row().cells
                for i, value in enumerate(row):
                    cells[i].text = str(value)
    for section in document.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)
    ensure_dir(output)
    document.save(output)


def render_xlsx_ooxml(spec, output):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)
    sheets = spec.get('sheets', []) or [{'name': spec.get('title', 'Sheet1'), 'headers': [], 'rows': []}]
    for sheet_spec in sheets:
        sheet = workbook.create_sheet(title=str(sheet_spec.get('name', 'Sheet'))[:31])
        row_index = 1
        headers = sheet_spec.get('headers', [])
        if headers:
            for col_index, value in enumerate(headers, start=1):
                cell = sheet.cell(row=row_index, column=col_index, value=value)
                cell.font = Font(bold=True, color='FFFFFF')
                cell.fill = PatternFill('solid', fgColor='1F4E79')
            row_index += 1
        for row in sheet_spec.get('rows', []):
            for col_index, value in enumerate(row, start=1):
                sheet.cell(row=row_index, column=col_index, value=value)
            row_index += 1
        max_col = max(len(headers), *(len(row) for row in sheet_spec.get('rows', [])), 1)
        for col_index in range(1, max_col + 1):
            sheet.column_dimensions[get_column_letter(col_index)].width = 18
        if headers:
            sheet.freeze_panes = 'A2'
    ensure_dir(output)
    workbook.save(output)


def render_pptx_ooxml(spec, output):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = spec.get('title', 'Slide Deck')
    slide.placeholders[1].text = f"Generated: {datetime.now():%Y-%m-%d %H:%M}"

    bullet_layout = prs.slide_layouts[1]
    for slide_spec in spec.get('slides', []):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = str(slide_spec.get('title', 'Untitled'))
        body = slide.placeholders[1].text_frame
        body.clear()
        subtitle = slide_spec.get('subtitle')
        if subtitle:
            p = body.paragraphs[0]
            p.text = str(subtitle)
            p.font.size = Pt(18)
        body_text = slide_spec.get('body')
        if body_text:
            p = body.add_paragraph() if subtitle else body.paragraphs[0]
            p.text = str(body_text)
            p.font.size = Pt(20)
        for item in slide_spec.get('bullets', []):
            p = body.add_paragraph()
            p.text = str(item)
            p.level = 0
            p.font.size = Pt(20)
        notes = slide_spec.get('notes')
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
    ensure_dir(output)
    prs.save(output)


def has_ooxml_deps(kind):
    modules = {'docx': 'docx', 'xlsx': 'openpyxl', 'pptx': 'pptx'}
    try:
        __import__(modules[kind])
        return True
    except Exception:
        return False


def main():
    parser = argparse.ArgumentParser(description='Generate Office documents from JSON specs.')
    parser.add_argument('kind', choices=['docx', 'xlsx', 'pptx'])
    parser.add_argument('spec')
    parser.add_argument('output')
    parser.add_argument('--format', choices=['auto', 'ooxml', 'html'], default='auto')
    args = parser.parse_args()
    spec = read_spec(args.spec)
    html_renderers = {'docx': render_docx_html, 'xlsx': render_xlsx_html, 'pptx': render_pptx_html}
    ooxml_renderers = {'docx': render_docx_ooxml, 'xlsx': render_xlsx_ooxml, 'pptx': render_pptx_ooxml}
    use_ooxml = args.format == 'ooxml' or (args.format == 'auto' and has_ooxml_deps(args.kind))
    if use_ooxml:
        ooxml_renderers[args.kind](spec, args.output)
    else:
        html_renderers[args.kind](spec, args.output)
    print(args.output)


if __name__ == '__main__':
    main()
