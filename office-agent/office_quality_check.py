#!/usr/bin/env python3
"""Structured quality checks for Office Agent outputs."""

import argparse
import os
import sys
import zipfile
from datetime import datetime

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
if SCRIPT_DIR not in sys.path:
    sys.path.insert(0, SCRIPT_DIR)

from office_common import read_json, write_json, write_text


def add_issue(result, section, level, code, message):
    issue = {'level': level, 'code': code, 'message': message}
    result[section]['issues'].append(issue)
    if level == 'fail':
        result['status'] = 'fail'
    elif level == 'warn' and result['status'] == 'pass':
        result['status'] = 'warn'


def extracted_has_attention_or_source(extracted):
    for slide in extracted.get('slides', []):
        for item in slide.get('items', []):
            text = item.get('text', '')
            if item.get('semantic_guess') in {'risk', 'source_or_disclaimer', 'summary_or_recommendation'}:
                return True
            if any(key in text.lower() for key in ['risk', 'source', 'disclaimer']):
                return True
            if any(key in text for key in ['风险', '来源', '免责声明', '建议']):
                return True
    return False


def read_docx_stats(path):
    from docx import Document

    doc = Document(path)
    text = '\n'.join(p.text for p in doc.paragraphs)
    return {
        'kind': 'docx',
        'paragraphs': len(doc.paragraphs),
        'tables': len(doc.tables),
        'text': text,
        'title': doc.paragraphs[0].text.strip() if doc.paragraphs else '',
    }


def read_pptx_stats(path):
    from pptx import Presentation

    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                texts.append(shape.text)
    return {
        'kind': 'pptx',
        'slides': len(prs.slides),
        'text_items': len(texts),
        'text': '\n'.join(texts),
        'title': texts[0].strip() if texts else '',
    }


def read_xlsx_stats(path):
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=False)
    text_values = []
    non_empty = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    non_empty += 1
                    text_values.append(str(cell.value))
    return {
        'kind': 'xlsx',
        'sheets': len(wb.worksheets),
        'sheet_names': wb.sheetnames,
        'non_empty_cells': non_empty,
        'text': '\n'.join(text_values),
        'title': wb.sheetnames[0] if wb.sheetnames else '',
    }


def read_stats(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.docx':
        return read_docx_stats(path)
    if ext == '.pptx':
        return read_pptx_stats(path)
    if ext == '.xlsx':
        return read_xlsx_stats(path)
    raise ValueError(f'Unsupported file extension for quality check: {ext}')


def check(output_path, extract_json=None, ledger_json=None, plan_json=None, change_log_json=None, risk_level=None):
    result = {
        'schema_version': '1.0',
        'checked_at': datetime.now().isoformat(timespec='seconds'),
        'file': output_path,
        'status': 'pass',
        'technical': {'checks': {}, 'issues': []},
        'content': {'checks': {}, 'issues': []},
        'experience': {'checks': {}, 'issues': []},
        'manual_review': {'items': [], 'issues': []},
        'warnings': [],
    }

    exists = os.path.exists(output_path)
    result['technical']['checks']['exists'] = exists
    if not exists:
        add_issue(result, 'technical', 'fail', 'missing_output', 'Output file does not exist.')
        result['warnings'] = [i['message'] for s in ['technical', 'content', 'experience'] for i in result[s]['issues'] if i['level'] == 'warn']
        return result

    size = os.path.getsize(output_path)
    result['technical']['checks']['size_bytes'] = size
    if size <= 0:
        add_issue(result, 'technical', 'fail', 'empty_output', 'Output file is empty.')
    if not zipfile.is_zipfile(output_path):
        add_issue(result, 'technical', 'fail', 'invalid_ooxml', 'Output is not a valid OOXML zip package.')

    stats = {}
    try:
        stats = read_stats(output_path)
        result['technical']['checks'].update({k: v for k, v in stats.items() if k != 'text'})
    except Exception as exc:
        add_issue(result, 'technical', 'fail', 'open_failed', f'Output could not be opened by parser: {exc}')

    plan = read_json(plan_json) if plan_json else None
    extracted = read_json(extract_json) if extract_json else None
    ledger = read_json(ledger_json) if ledger_json else None
    change_log = read_json(change_log_json) if change_log_json else None
    effective_risk = risk_level or (plan or {}).get('risk_level') or 'unknown'
    result['content']['checks']['risk_level'] = effective_risk

    text = stats.get('text', '')
    has_title = bool(stats.get('title'))
    result['experience']['checks']['has_cover_like_title'] = has_title
    if not has_title:
        add_issue(result, 'experience', 'warn', 'missing_title', 'Output does not appear to start with a title.')

    has_summary = any(key in text for key in ['执行摘要', 'Executive Summary', '内容导航', '一页摘要', 'Summary'])
    result['experience']['checks']['has_summary_or_navigation'] = has_summary
    if stats.get('kind') == 'docx' and not has_summary and effective_risk in {'medium', 'high'}:
        add_issue(result, 'experience', 'warn', 'missing_summary_navigation', 'Medium/high-risk DOCX output has no obvious summary or navigation section.')

    if extracted:
        slides = len(extracted.get('slides', []))
        result['content']['checks']['source_slides'] = slides
        result['content']['checks']['has_source_map'] = True
    else:
        result['content']['checks']['has_source_map'] = False

    if ledger:
        rows = ledger.get('rows', [])
        result['content']['checks']['ledger_rows'] = len(rows)
        if extracted and len(rows) != result['content']['checks'].get('source_slides'):
            add_issue(result, 'content', 'warn', 'ledger_slide_mismatch', 'Ledger row count does not match source slide count.')
    elif effective_risk in {'medium', 'high'}:
        add_issue(result, 'content', 'fail', 'missing_ledger', 'Medium/high-risk material must include a fidelity ledger.')

    if change_log:
        result['content']['checks']['change_log_entries'] = len(change_log.get('changes', []))
    elif plan and plan.get('task_type') == 'modify':
        add_issue(result, 'content', 'fail', 'missing_change_log', 'Modify tasks must include a change log.')

    has_attention = any(key in text for key in ['风险', 'risk', 'Risk', '来源', 'source', 'Source', '免责声明', '建议'])
    if not has_attention and extracted:
        has_attention = extracted_has_attention_or_source(extracted)
    result['experience']['checks']['has_attention_or_source_content'] = has_attention
    if effective_risk in {'medium', 'high'} and not has_attention:
        add_issue(result, 'experience', 'warn', 'missing_risk_or_source_cue', 'Risk/source/disclaimer cues were not found; manual review is recommended.')

    result['experience']['checks']['avoids_slide_by_slide'] = 'Slide 1:' not in text[:2000]
    if not result['experience']['checks']['avoids_slide_by_slide'] and effective_risk in {'medium', 'high'}:
        add_issue(result, 'experience', 'warn', 'slide_by_slide_style', 'Medium/high-risk output still appears to be slide-by-slide rather than report-like.')

    if effective_risk in {'medium', 'high'}:
        result['manual_review']['items'].append('Verify professional facts, numbers, risks, and source/disclaimer handling manually.')
    if stats.get('kind') in {'pptx', 'docx'}:
        result['manual_review']['items'].append('Review visual layout and business wording manually; these are not fully deterministic checks.')

    result['warnings'] = [i['message'] for s in ['technical', 'content', 'experience'] for i in result[s]['issues'] if i['level'] == 'warn']
    return result


def markdown(result):
    lines = [
        '# Office Quality Report',
        '',
        f"- File: `{result.get('file')}`",
        f"- Status: `{result.get('status')}`",
        f"- Checked at: `{result.get('checked_at')}`",
        '',
    ]
    for section in ['technical', 'content', 'experience']:
        lines.append(f"## {section.title()}")
        lines.append('')
        checks = result.get(section, {}).get('checks', {})
        for key, value in checks.items():
            lines.append(f'- {key}: `{value}`')
        issues = result.get(section, {}).get('issues', [])
        if issues:
            lines.append('')
            lines.append('Issues:')
            for issue in issues:
                lines.append(f"- `{issue['level']}` `{issue['code']}`: {issue['message']}")
        lines.append('')
    lines.append('## Manual Review')
    lines.append('')
    items = result.get('manual_review', {}).get('items', [])
    if items:
        for item in items:
            lines.append(f'- {item}')
    else:
        lines.append('- None')
    lines.append('')
    return '\n'.join(lines)


def main():
    parser = argparse.ArgumentParser(description='Run structured quality checks for Office Agent output.')
    parser.add_argument('output_file')
    parser.add_argument('output_json')
    parser.add_argument('--output-md', default='')
    parser.add_argument('--extract-json', default='')
    parser.add_argument('--ledger-json', default='')
    parser.add_argument('--plan-json', default='')
    parser.add_argument('--change-log-json', default='')
    parser.add_argument('--risk-level', default='')
    args = parser.parse_args()

    result = check(
        args.output_file,
        extract_json=args.extract_json or None,
        ledger_json=args.ledger_json or None,
        plan_json=args.plan_json or None,
        change_log_json=args.change_log_json or None,
        risk_level=args.risk_level or None,
    )
    write_json(args.output_json, result)
    if args.output_md:
        write_text(args.output_md, markdown(result))
    print(args.output_json)


if __name__ == '__main__':
    main()
