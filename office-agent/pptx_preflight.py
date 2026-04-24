#!/usr/bin/env python3
"""Preflight a PPTX file before conversion."""

import argparse
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
if SCRIPT_DIR not in sys.path:
    sys.path.insert(0, SCRIPT_DIR)

from office_common import classify_text_density, detect_document_type, risk_level_for, write_json
from template_registry import recommended_modes


def shape_text(shape):
    if getattr(shape, 'has_text_frame', False):
        return shape.text.strip()
    return ''


def has_notes(slide):
    try:
        return bool(slide.notes_slide.notes_text_frame.text.strip())
    except Exception:
        return False


def preflight(input_path):
    prs = Presentation(input_path)
    titles = []
    texts = []
    tables = 0
    pictures = 0
    charts = 0
    groups = 0
    notes = 0
    smartart_like = 0

    for slide in prs.slides:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            titles.append(slide.shapes.title.text.strip())
        if has_notes(slide):
            notes += 1
        for shape in slide.shapes:
            text = shape_text(shape)
            if text:
                texts.extend([line.strip() for line in text.splitlines() if line.strip()])
            if getattr(shape, 'has_table', False):
                tables += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                charts += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                groups += 1
            if shape.shape_type in {MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.FREEFORM}:
                smartart_like += 1

    doc_type = detect_document_type(titles, texts)
    risk = risk_level_for(doc_type)
    density = classify_text_density(len(prs.slides), len(texts))
    result = {
        'file': input_path,
        'file_size_bytes': os.path.getsize(input_path),
        'slides': len(prs.slides),
        'text_items': len(texts),
        'text_density': density,
        'tables': tables,
        'images': pictures,
        'charts': charts,
        'group_shapes': groups,
        'smartart_like_shapes': smartart_like,
        'slides_with_speaker_notes': notes,
        'detected_type': doc_type,
        'risk_level': risk,
        'requires_confirmation': risk in {'medium', 'high'},
        'warnings': build_warnings(risk, charts, smartart_like, pictures),
    }
    result['recommended_modes'] = recommended_modes(result)
    return result


def build_warnings(risk, charts, smartart_like, pictures):
    warnings = []
    if risk == 'high':
        warnings.append('High-risk document type: confirm fidelity level before professional rewriting.')
    if charts:
        warnings.append('Charts may be preserved as text/image summaries; verify manually.')
    if smartart_like:
        warnings.append('SmartArt/group shapes may not be fully structured; ledger should mark handling.')
    if pictures:
        warnings.append('Images can be embedded or omitted depending on selected conversion plan.')
    return warnings


def main():
    parser = argparse.ArgumentParser(description='Preflight a PPTX file before Office Agent conversion.')
    parser.add_argument('input')
    parser.add_argument('output')
    args = parser.parse_args()
    result = preflight(args.input)
    write_json(args.output, result)
    print(args.output)


if __name__ == '__main__':
    main()
