#!/usr/bin/env python3
"""Create small PPTX regression fixtures for Office Agent tests."""

import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        p = body.paragraphs[0] if index == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide


def text_only(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Text Only Project Brief'
    slide.placeholders[1].text = 'Summary and next steps'
    add_bullet_slide(prs, 'Executive Summary', ['Goal: validate conversion.', 'Risk: low.', 'Source: synthetic fixture.'])
    add_bullet_slide(prs, 'Next Steps', ['Generate plan.', 'Run conversion.', 'Check quality.'])
    prs.save(path)


def visual_mix(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Visual Mix Business Update'
    slide.placeholders[1].text = 'Contains text, table, and chart-like content'
    slide = add_bullet_slide(prs, 'Market Snapshot', ['Business momentum is improving.', 'Source: synthetic fixture.', 'Risk: medium review needed.'])
    rows, cols = 3, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(3.5), Inches(5.8), Inches(1.2)).table
    for c, value in enumerate(['Metric', 'Current', 'Target']):
        table.cell(0, c).text = value
    for r, values in enumerate([['Revenue', '100', '120'], ['Cost', '60', '55']], start=1):
        for c, value in enumerate(values):
            table.cell(r, c).text = value
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3']
    chart_data.add_series('Revenue', (10, 15, 20))
    prs.slides.add_slide(prs.slide_layouts[5]).shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(6), Inches(4), chart_data
    )
    prs.save(path)


def high_density(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'Investment Review High Density Fixture'
    slide.placeholders[1].text = 'High-risk synthetic investment material'
    for idx in range(1, 5):
        bullets = [
            f'Investment thesis item {idx}.{j}: market, product, finance, risk, source, disclaimer.'
            for j in range(1, 18)
        ]
        add_bullet_slide(prs, f'Investment Risk Assessment {idx}', bullets)
    prs.save(path)


def main():
    root = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'examples', 'regression_inputs')
    ensure_dir(root)
    text_only(os.path.join(root, 'text_only.pptx'))
    visual_mix(os.path.join(root, 'visual_mix.pptx'))
    high_density(os.path.join(root, 'high_density.pptx'))
    print(root)


if __name__ == '__main__':
    main()
